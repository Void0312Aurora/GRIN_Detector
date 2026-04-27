from __future__ import annotations

import math
import shutil
import tempfile
from pathlib import Path

import matplotlib.pyplot as plt
import fitz
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from spire.presentation import FileFormat, Presentation as SpirePresentation


ROOT = Path(__file__).resolve().parents[1]
REPORTS_DIR = ROOT / "docs" / "reports"
ASSETS_DIR = REPORTS_DIR / "答辩汇报_2026-04-23_assets"
OUTPUT_PPTX = REPORTS_DIR / "答辩汇报_2026-04-23.pptx"
OUTPUT_PDF = REPORTS_DIR / "答辩汇报_2026-04-23.pdf"
PREVIEW_DIR = REPORTS_DIR / "答辩汇报_2026-04-23_preview"
PREVIEW_OVERVIEW = REPORTS_DIR / "答辩汇报_2026-04-23_preview_overview.png"

IMG_HERO = (
    REPORTS_DIR
    / "中期检查附件_2026-04-08/02_成熟度附件/05_拿点数据_三组数据与说明/05_图1_代表性预测结果.png"
)
IMG_SUCCESS = (
    REPORTS_DIR
    / "中期检查附件_2026-04-08/02_成熟度附件/06_分析结果_图表与讨论/06_代表性成功样例.png"
)
IMG_FAILURE = (
    REPORTS_DIR
    / "中期检查附件_2026-04-08/02_成熟度附件/06_分析结果_图表与讨论/06_传统基线失败样例.png"
)

FONT_CN = "Microsoft YaHei"
FONT_EN = "Aptos"
PAGE_W = 13.333
PAGE_H = 7.5
LEFT = 0.68
RIGHT = 12.66
CONTENT_TOP = 1.48
CONTENT_BOTTOM = 6.82

COLORS = {
    "navy": "0F172A",
    "blue": "1D4ED8",
    "sky": "0284C7",
    "teal": "0F766E",
    "green": "059669",
    "orange": "D97706",
    "red": "DC2626",
    "ink": "1E293B",
    "muted": "475569",
    "line": "CBD5E1",
    "bg": "F8FAFC",
    "card": "FFFFFF",
    "soft_blue": "EAF2FF",
    "soft_green": "EAF9F1",
    "soft_red": "FEECEC",
    "soft_orange": "FFF5E8",
    "soft_ink": "E2E8F0",
}


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def fill_background(slide, color: str = COLORS["bg"]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def set_text_style(run, size: int, color: str, bold: bool = False, name: str = FONT_CN) -> None:
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: int = 18,
    color: str = COLORS["ink"],
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    name: str = FONT_CN,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    set_text_style(run, size=size, color=color, bold=bold, name=name)
    return box


def add_lines(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: list[str],
    *,
    size: int = 16,
    color: str = COLORS["ink"],
    bullet: str | None = None,
    line_spacing: float = 1.18,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = f"{bullet} {line}" if bullet else line
        set_text_style(run, size=size, color=color, bold=False)
    return box


def add_panel(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    title: str,
    body: str | None = None,
    accent: str = COLORS["blue"],
    fill: str = COLORS["card"],
    title_size: int = 18,
    body_size: int = 14,
    title_box_height: float = 0.32,
    body_top: float = 0.54,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(COLORS["line"])
    shape.line.width = Pt(1.0)

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(0.12),
        Inches(height),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(accent)
    bar.line.fill.background()

    add_textbox(
        slide,
        left + 0.24,
        top + 0.16,
        width - 0.36,
        title_box_height,
        title,
        size=title_size,
        color=COLORS["navy"],
        bold=True,
    )
    if body:
        add_textbox(
            slide,
            left + 0.24,
            top + body_top,
            width - 0.36,
            max(height - body_top - 0.14, 0.25),
            body,
            size=body_size,
            color=COLORS["muted"],
        )
    return shape


def add_stat_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    label: str,
    value: str,
    detail: str,
    accent: str,
    value_size: int | None = None,
    detail_size: int = 11,
):
    value_size = value_size or (18 if len(value) > 12 else 23)
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS["card"])
    shape.line.color.rgb = rgb(COLORS["line"])
    shape.line.width = Pt(1.0)

    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left + 0.12),
        Inches(top + 0.12),
        Inches(width - 0.24),
        Inches(0.24),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = rgb(accent)
    tag.fill.transparency = 0.88
    tag.line.fill.background()

    add_textbox(slide, left + 0.22, top + 0.10, width - 0.44, 0.18, label, size=10, color=accent, bold=True, name=FONT_EN)
    add_textbox(
        slide,
        left + 0.18,
        top + 0.44,
        width - 0.36,
        0.34,
        value,
        size=value_size,
        color=COLORS["navy"],
        bold=True,
        name=FONT_EN,
    )
    add_textbox(
        slide,
        left + 0.18,
        top + 0.80,
        width - 0.36,
        max(height - 0.92, 0.22),
        detail,
        size=detail_size,
        color=COLORS["muted"],
    )
    return shape


def add_header(slide, title: str, kicker: str, page_number: int) -> None:
    add_textbox(slide, LEFT, 0.34, 2.6, 0.20, kicker, size=10, color=COLORS["blue"], bold=True, name=FONT_EN)
    add_textbox(slide, LEFT, 0.60, 8.8, 0.42, title, size=27, color=COLORS["navy"], bold=True)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(LEFT),
        Inches(1.18),
        Inches(RIGHT - LEFT),
        Inches(0.03),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(COLORS["soft_ink"])
    line.line.fill.background()
    add_textbox(
        slide,
        12.00,
        0.40,
        0.62,
        0.22,
        f"{page_number:02d}",
        size=11,
        color=COLORS["muted"],
        bold=True,
        align=PP_ALIGN.RIGHT,
        name=FONT_EN,
    )


def add_footer(slide, text: str) -> None:
    add_textbox(slide, LEFT, 7.08, 8.4, 0.16, text, size=9, color=COLORS["muted"])


def crop_to_ratio(src: Path, dst: Path, ratio: float) -> None:
    with Image.open(src) as img:
        w, h = img.size
        current = w / h
        if current > ratio:
            new_w = int(h * ratio)
            left = (w - new_w) // 2
            cropped = img.crop((left, 0, left + new_w, h))
        else:
            new_h = int(w / ratio)
            top = (h - new_h) // 2
            cropped = img.crop((0, top, w, top + new_h))
        cropped.save(dst)


def build_chart_ablation(dst: Path) -> None:
    methods = [
        "Poisson",
        "Oracle",
        "NN default",
        "Teacher",
        "Prior",
        "Prior + Input",
        "Prior + Teacher",
    ]
    f1 = [0.0507, 0.0515, 0.0013, 0.0309, 0.1892, 0.2427, 0.0509]
    rmse = [0.0625, 0.0614, 0.0153, 0.0591, 0.0117, 0.0100, 0.0614]
    colors = ["#E2E8F0"] * len(methods)
    colors[0] = "#FDE7C7"
    colors[5] = f"#{COLORS['blue']}"

    fig, axes = plt.subplots(1, 2, figsize=(11.6, 4.7), dpi=220)
    axes[0].barh(methods, f1, color=colors, edgecolor="#CBD5E1")
    axes[0].invert_yaxis()
    axes[0].set_xlim(0, 0.28)
    axes[0].grid(axis="x", alpha=0.18)
    axes[0].set_title("Defect F1", fontsize=14, fontweight="bold")
    axes[0].spines[["top", "right", "left"]].set_visible(False)
    for idx, value in enumerate(f1):
        axes[0].text(value + 0.004, idx, f"{value:.4f}", va="center", fontsize=9)

    axes[1].barh(methods, rmse, color=colors, edgecolor="#CBD5E1")
    axes[1].invert_yaxis()
    axes[1].set_xlim(0, 0.07)
    axes[1].grid(axis="x", alpha=0.18)
    axes[1].set_title("Global RMSE", fontsize=14, fontweight="bold")
    axes[1].spines[["top", "right", "left"]].set_visible(False)
    for idx, value in enumerate(rmse):
        axes[1].text(value + 0.001, idx, f"{value:.4f}", va="center", fontsize=9)

    fig.suptitle("Ablation evidence: residual-prior input is the best current path", fontsize=16, fontweight="bold")
    fig.tight_layout(rect=(0, 0, 1, 0.94))
    fig.savefig(dst, bbox_inches="tight")
    plt.close(fig)


def build_chart_ood(dst: Path) -> None:
    scenarios = ["ID", "Scratch", "Dot", "Small", "Large", "Pos shift"]
    poisson = [0.7515, 0.7315, 0.7142, 0.4129, 0.8274, 0.8811]
    nn = [0.7638, 0.7543, 0.7460, 0.2260, 0.8769, 0.8738]
    x = range(len(scenarios))

    fig, ax = plt.subplots(figsize=(10.8, 4.8), dpi=220)
    width = 0.36
    ax.bar([i - width / 2 for i in x], poisson, width=width, label="Poisson", color=f"#{COLORS['orange']}", alpha=0.72)
    ax.bar([i + width / 2 for i in x], nn, width=width, label="NN", color=f"#{COLORS['blue']}")
    ax.set_ylim(0, 1.0)
    ax.set_ylabel("Defect F1")
    ax.set_xticks(list(x), scenarios)
    ax.grid(axis="y", alpha=0.18)
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_title("OOD scenarios: strong on common defects, weak on very small defects", fontsize=15, fontweight="bold")
    ax.legend(frameon=False, loc="upper right")
    fig.tight_layout()
    fig.savefig(dst, bbox_inches="tight")
    plt.close(fig)


def build_chart_qc(dst: Path) -> None:
    methods = ["Prior + Input", "Reflect", "Oracle", "Poisson"]
    pass_rate = [1.0, 0.975, 0.0, 0.0]
    colors = [f"#{COLORS['blue']}", f"#{COLORS['sky']}", "#D8DEE7", "#D8DEE7"]

    fig, ax = plt.subplots(figsize=(9.4, 4.6), dpi=220)
    bars = ax.bar(methods, pass_rate, color=colors, edgecolor="#CBD5E1")
    ax.set_ylim(0, 1.1)
    ax.set_ylabel("QC pass rate")
    ax.grid(axis="y", alpha=0.18)
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_title("QC gating separates reliable neural results from unstable baselines", fontsize=15, fontweight="bold")
    for bar, value in zip(bars, pass_rate):
        ax.text(bar.get_x() + bar.get_width() / 2, value + 0.03, f"{value:.1%}", ha="center", fontsize=10)
    fig.tight_layout()
    fig.savefig(dst, bbox_inches="tight")
    plt.close(fig)


def build_chart_pre_real(dst: Path) -> None:
    fig, axes = plt.subplots(2, 2, figsize=(11.0, 5.7), dpi=220)

    axes[0, 0].plot(["64", "256", "512"], [0.925, 0.0, 0.0], marker="o", color=f"#{COLORS['blue']}", linewidth=2.5)
    axes[0, 0].set_title("Resolution")
    axes[0, 1].plot(["0.002", "0.005", "0.010", "0.020"], [0.0, 0.925, 0.0, 0.0], marker="o", color=f"#{COLORS['red']}", linewidth=2.5)
    axes[0, 1].set_title("Noise")
    axes[1, 0].plot(["0.95", "0.98", "1.00", "1.02", "1.05"], [1.0, 0.975, 0.925, 0.900, 0.875], marker="o", color=f"#{COLORS['green']}", linewidth=2.5)
    axes[1, 0].set_title("Aperture radius")
    axes[1, 1].plot(["0.95", "0.98", "1.00", "1.02", "1.05"], [0.900, 0.900, 0.925, 0.950, 0.900], marker="o", color=f"#{COLORS['teal']}", linewidth=2.5)
    axes[1, 1].set_title("dx scale")

    for ax in axes.flatten():
        ax.set_ylim(0, 1.05)
        ax.set_ylabel("QC pass rate")
        ax.grid(alpha=0.18)
        ax.spines[["top", "right"]].set_visible(False)

    fig.suptitle("Pre-real validation: stable under small calibration bias, sensitive to strong domain shift", fontsize=15, fontweight="bold")
    fig.tight_layout(rect=(0, 0, 1, 0.94))
    fig.savefig(dst, bbox_inches="tight")
    plt.close(fig)


def build_assets() -> dict[str, Path]:
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    assets = {
        "cover": ASSETS_DIR / "cover_cropped.png",
        "success": ASSETS_DIR / "success_cropped.png",
        "failure": ASSETS_DIR / "failure_cropped.png",
        "ablation": ASSETS_DIR / "ablation_chart.png",
        "ood": ASSETS_DIR / "ood_chart.png",
        "qc": ASSETS_DIR / "qc_chart.png",
        "pre_real": ASSETS_DIR / "pre_real_chart.png",
    }

    crop_to_ratio(IMG_HERO, assets["cover"], 16 / 9)
    crop_to_ratio(IMG_SUCCESS, assets["success"], 1.48)
    crop_to_ratio(IMG_FAILURE, assets["failure"], 1.48)
    build_chart_ablation(assets["ablation"])
    build_chart_ood(assets["ood"])
    build_chart_qc(assets["qc"])
    build_chart_pre_real(assets["pre_real"])
    return assets


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(PAGE_W)
    prs.slide_height = Inches(PAGE_H)
    return prs


def add_cover_slide(prs: Presentation, assets: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(assets["cover"]), 0, 0, width=prs.slide_width, height=prs.slide_height)

    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = rgb(COLORS["navy"])
    overlay.fill.transparency = 0.28
    overlay.line.fill.background()

    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.76), Inches(0.55), Inches(3.65), Inches(0.34))
    pill.fill.solid()
    pill.fill.fore_color.rgb = rgb(COLORS["blue"])
    pill.fill.transparency = 0.10
    pill.line.fill.background()
    add_textbox(slide, 0.95, 0.61, 3.20, 0.16, "DESIGN DEFENSE DRAFT | 2026-04-23", size=10, color="FFFFFF", bold=True, name=FONT_EN)

    add_textbox(
        slide,
        0.82,
        1.25,
        7.7,
        1.30,
        "基于物理信息神经网络和液晶微分的\n微透镜缺陷检测系统",
        size=27,
        color="FFFFFF",
        bold=True,
    )
    add_textbox(
        slide,
        0.86,
        2.88,
        6.2,
        0.70,
        "答辩汇报草案\n聚焦当前已有成果、证据链和下一步落地路径",
        size=17,
        color="E2E8F0",
    )

    add_panel(
        slide,
        0.84,
        5.48,
        2.32,
        1.18,
        title="阶段定位",
        body="成熟度 6\n定位在分析结果阶段",
        accent=COLORS["sky"],
        fill=COLORS["card"],
        title_size=15,
        body_size=11,
        title_box_height=0.24,
        body_top=0.48,
    )
    add_panel(
        slide,
        3.38,
        5.48,
        2.98,
        1.18,
        title="证据链",
        body="消融 / OOD / QC\n证据链已闭环",
        accent=COLORS["blue"],
        fill=COLORS["card"],
        title_size=15,
        body_size=11,
        title_box_height=0.24,
        body_top=0.48,
    )
    add_panel(
        slide,
        6.58,
        5.48,
        3.04,
        1.18,
        title="下一步",
        body="真实 pilot 闭环\n面向真实闭环与成果凝练",
        accent=COLORS["green"],
        fill=COLORS["card"],
        title_size=15,
        body_size=11,
        title_box_height=0.24,
        body_top=0.48,
    )
    add_textbox(slide, 9.72, 7.02, 2.90, 0.16, "Source: 中期检查报告与成熟度附件", size=9, color="E2E8F0", align=PP_ALIGN.RIGHT, name=FONT_EN)


def add_problem_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide)
    add_header(slide, "1. 研究背景与核心问题", "WHY THIS PROJECT", 2)
    add_textbox(slide, LEFT, 1.42, 9.4, 0.42, "目标不是单纯提高重建精度，而是在更简化硬件条件下稳定识别微透镜缺陷，并判断结果是否可信。", size=20, color=COLORS["navy"], bold=True)

    add_panel(slide, 0.74, 2.02, 4.15, 1.22, title="挑战 1 | 无偏置条件下符号难恢复", body="差分强度有响应，但相位符号恢复困难，传统流程容易丢失关键缺陷信息。", accent=COLORS["red"])
    add_panel(slide, 0.74, 3.45, 4.15, 1.22, title="挑战 2 | 积分伪影影响可解释性", body="Poisson 类基线在边缘、孔径外和局部异常区更容易出现伪影。", accent=COLORS["orange"])
    add_panel(slide, 0.74, 4.88, 4.15, 1.18, title="挑战 3 | 真实标定与弱真值获取难", body="真实实验既缺直接真值，也涉及跨模态配准和标定扰动问题。", accent=COLORS["teal"])

    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.95), Inches(7.0), Inches(4.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS["navy"])
    shape.line.fill.background()
    add_textbox(slide, 5.52, 2.20, 2.2, 0.18, "本项目的回答方式", size=12, color="93C5FD", bold=True)
    add_textbox(slide, 5.52, 2.52, 5.8, 0.54, "差分物理模型 + 物理信息神经网络 + 质量门控", size=24, color="FFFFFF", bold=True)
    add_lines(
        slide,
        5.55,
        3.28,
        5.95,
        1.48,
        [
            "用物理一致性约束替代单纯积分或单纯拟合。",
            "用标准面与稀疏残差先验提高缺陷恢复能力。",
            "用 QC 门控判断一次重建是否足够可信。",
        ],
        size=17,
        color="E2E8F0",
        bullet="-",
    )
    add_panel(slide, 5.55, 5.10, 1.92, 1.08, title="问题", body="先判断路线可行性", accent=COLORS["sky"], fill=COLORS["card"], title_size=14, body_size=10, title_box_height=0.24, body_top=0.46)
    add_panel(slide, 7.68, 5.10, 1.92, 1.08, title="状态", body="典型缺陷场景已跑通", accent=COLORS["green"], fill=COLORS["card"], title_size=14, body_size=10, title_box_height=0.24, body_top=0.46)
    add_panel(slide, 9.81, 5.10, 1.92, 1.08, title="边界", body="超小缺陷仍然困难", accent=COLORS["red"], fill=COLORS["card"], title_size=14, body_size=10, title_box_height=0.24, body_top=0.46)
    add_footer(slide, "问题定义和项目摘要来自中期检查报告。")


def add_route_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide)
    add_header(slide, "2. 技术路线与证据链设计", "TECHNICAL ROUTE", 3)
    add_textbox(slide, LEFT, 1.42, 10.3, 0.38, "答辩里最关键的不是“做了哪些实验”，而是“这些实验如何共同支持当前路线成立”。", size=20, color=COLORS["navy"], bold=True)

    steps = [
        ("差分输入", "无偏置差分观测"),
        ("物理重建", "前向一致性进入主线"),
        ("残差先验", "标准面 + 稀疏残差"),
        ("质量门控", "判断结果是否可信"),
    ]
    x_positions = [0.78, 3.72, 6.66, 9.60]
    for idx, ((title, body), left) in enumerate(zip(steps, x_positions), 1):
        add_panel(slide, left, 2.08, 2.35, 1.42, title=f"{idx}. {title}", body=body, accent=COLORS["blue"] if idx < 3 else COLORS["green"], title_size=16, body_size=12)
        if idx < len(steps):
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(left + 2.15), Inches(2.50), Inches(0.34), Inches(0.42))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = rgb(COLORS["line"])
            arrow.line.fill.background()

    add_panel(
        slide,
        0.78,
        3.90,
        11.4,
        1.05,
        title="当前证据链",
        body="消融回答“为什么当前方案成立”；OOD 回答“泛化到哪里”；QC 回答“结果是否可信”；预真实验证回答“离真实应用还有多远”。",
        accent=COLORS["navy"],
        fill=COLORS["soft_blue"],
        body_size=14,
    )

    add_panel(slide, 0.78, 5.28, 5.35, 1.22, title="已经完成", body="问题聚焦、算法原型、关键消融、OOD、QC、预真实验证，以及 qDIC / WLI 方案准备。", accent=COLORS["blue"], fill=COLORS["card"])
    add_panel(slide, 6.45, 5.28, 5.73, 1.22, title="仍待推进", body="正式真实闭环、真实样本阈值重校、论文初稿和专利文本形成。", accent=COLORS["red"], fill=COLORS["card"])
    add_footer(slide, "证据链结构来自中期报告与成熟度 5/6 附件。")


def add_method_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide)
    add_header(slide, "3. 当前方法的核心设计", "METHOD DESIGN", 4)
    add_textbox(slide, LEFT, 1.42, 11.2, 0.40, "这条路线的价值不在于“换了一个网络”，而在于把物理、残差先验和可信度筛查放进同一个工作流。", size=20, color=COLORS["navy"], bold=True)

    add_panel(
        slide,
        0.78,
        2.05,
        3.7,
        1.95,
        title="A. Physics-informed\nreconstruction",
        body="用物理一致性项、标准 / 参考差分约束和稳定性项约束重建，避免纯 Poisson 与纯学习各自的短板。",
        accent=COLORS["blue"],
        title_size=13,
        body_size=12,
        title_box_height=0.54,
        body_top=0.86,
    )
    add_panel(
        slide,
        4.82,
        2.05,
        3.7,
        1.95,
        title="B. Residual-prior\nmodeling",
        body="把标准面与局部稀疏缺陷残差统一建模，现阶段优选实现是 residual-prior input。",
        accent=COLORS["teal"],
        title_size=13,
        body_size=12,
        title_box_height=0.54,
        body_top=0.86,
    )
    add_panel(
        slide,
        8.86,
        2.05,
        3.28,
        1.95,
        title="C. Quality-control\ngating",
        body="用边缘、视场外和物理残差等统计量筛查一次重建是否可信。",
        accent=COLORS["green"],
        title_size=13,
        body_size=12,
        title_box_height=0.54,
        body_top=0.86,
    )

    add_panel(slide, 0.78, 4.42, 11.36, 1.62, title="为什么不是传统基线", body="纯 Poisson 更容易在边缘、孔径外和局部异常区产生不可信伪影；纯数据驱动在无偏置符号恢复问题上又不够稳定。当前主线是在两者之间建立可解释的折中。", accent=COLORS["navy"], fill=COLORS["soft_blue"], body_size=15)
    add_panel(slide, 0.78, 6.10, 11.36, 0.94, title="一句话记忆点", body="当前方法的关键词不是“更复杂的网络”，而是“物理约束 + 残差先验 + 质量门控”。", accent=COLORS["blue"], fill=COLORS["card"], title_size=15, body_size=12, title_box_height=0.24, body_top=0.44)
    add_footer(slide, "方法设计依据中期报告“核心算法原型构建”和“质量门控方案”章节。")


def add_ablation_slide(prs: Presentation, assets: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide)
    add_header(slide, "4. 关键消融: 当前有效路径已经明确", "ABLATION", 5)

    slide.shapes.add_picture(str(assets["ablation"]), Inches(0.72), Inches(1.68), width=Inches(6.65))
    add_stat_box(slide, 7.75, 1.82, 2.12, 1.24, label="BEST F1", value="0.2427", detail="vs Poisson 0.0507", accent=COLORS["blue"], detail_size=10)
    add_stat_box(slide, 10.00, 1.82, 2.12, 1.24, label="BEST RMSE", value="0.0100", detail="baseline 0.0625", accent=COLORS["green"], detail_size=10)
    add_panel(
        slide,
        7.72,
        3.30,
        4.40,
        2.05,
        title="这页要传达的结论",
        body="1. 默认网络和仅教师项方案不足以稳定恢复缺陷。\n2. 引入残差先验后性能明显提升。\n3. residual-prior input 是当前阶段最优实现路径。",
        accent=COLORS["navy"],
        body_size=14,
    )
    add_panel(slide, 7.72, 5.62, 4.40, 0.96, title="一句话结论", body="当前路线不是“能不能跑”，而是已经出现了可重复、可解释的优选实现。", accent=COLORS["blue"], fill=COLORS["soft_blue"], title_size=15, body_size=12, title_box_height=0.24, body_top=0.44)
    add_footer(slide, "数据来自成熟度 5 数据组 1 和中期报告中的消融结论。")


def add_visual_slide(prs: Presentation, assets: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide)
    add_header(slide, "5. 代表性结果: 当前方法能恢复结构, 基线易失真", "VISUAL EVIDENCE", 6)

    add_panel(slide, 0.78, 1.64, 5.55, 4.70, title="OURS | 当前方法", accent=COLORS["blue"], title_size=16)
    add_panel(slide, 6.62, 1.64, 5.55, 4.70, title="BASELINE | 传统积分基线", accent=COLORS["orange"], title_size=16)
    slide.shapes.add_picture(str(assets["success"]), Inches(1.06), Inches(2.16), width=Inches(4.80), height=Inches(3.22))
    slide.shapes.add_picture(str(assets["failure"]), Inches(6.90), Inches(2.16), width=Inches(4.80), height=Inches(3.22))
    add_textbox(slide, 1.00, 5.84, 5.06, 0.38, "能够看到更清晰的缺陷响应和结构恢复，适合进入后续 QC 与分析链路。", size=12, color=COLORS["muted"])
    add_textbox(slide, 6.84, 5.84, 5.06, 0.38, "在边缘和局部区域更容易出现伪影或错误结构，可信度不足。", size=12, color=COLORS["muted"])

    add_panel(slide, 0.78, 6.30, 11.4, 0.48, title="这页的重点", body="当前方法已经能在代表性样例上恢复关键结构，而传统基线更容易在边缘和局部异常区失真。", accent=COLORS["navy"], fill=COLORS["soft_blue"], title_size=14, body_size=11)
    add_footer(slide, "图像来自成熟度 6 的代表性成功样例和传统基线失败样例。")


def add_ood_slide(prs: Presentation, assets: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide)
    add_header(slide, "6. OOD 泛化: 常见缺陷可用, 超小缺陷仍是短板", "OOD TEST", 7)

    slide.shapes.add_picture(str(assets["ood"]), Inches(0.72), Inches(1.70), width=Inches(6.85))
    add_stat_box(slide, 7.88, 1.92, 2.12, 1.18, label="LARGE DEFECT", value="0.8769", detail="大缺陷场景最强", accent=COLORS["green"], detail_size=10)
    add_stat_box(slide, 10.00, 1.92, 2.12, 1.18, label="SMALL DEFECT", value="0.2260", detail="超小缺陷明显掉点", accent=COLORS["red"], detail_size=10)
    add_panel(
        slide,
        7.86,
        3.30,
        4.24,
        2.15,
        title="答辩时建议这样讲",
        body="1. 常见划痕、点缺陷和大缺陷场景下，NN 方案已经具备竞争力。\n2. 这说明路线不是只在单一分布下有效。\n3. 但超小缺陷仍是当前最明确的瓶颈。",
        accent=COLORS["navy"],
        body_size=14,
    )
    add_panel(slide, 7.86, 5.72, 4.24, 0.98, title="风险态度", body="要主动承认边界，不要把所有场景都讲成已经解决。", accent=COLORS["red"], fill=COLORS["soft_red"], title_size=15, body_size=12, title_box_height=0.24, body_top=0.44)
    add_footer(slide, "数据来自成熟度 5 数据组 2 与成熟度 6 OOD 详细报告。")


def add_qc_slide(prs: Presentation, assets: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide)
    add_header(slide, "7. 质量门控: 不只出结果, 还判断结果是否可信", "QUALITY CONTROL", 8)

    slide.shapes.add_picture(str(assets["qc"]), Inches(0.72), Inches(1.74), width=Inches(6.35))
    add_stat_box(slide, 7.42, 1.86, 2.18, 1.14, label="PASS RATE", value="100%", detail="Prior + Input 全量通过", accent=COLORS["blue"], detail_size=10)
    add_stat_box(slide, 9.82, 1.86, 2.18, 1.14, label="BASELINE", value="0%", detail="Poisson / Oracle 全部失败", accent=COLORS["red"], detail_size=10)
    add_stat_box(slide, 7.42, 3.22, 2.18, 1.14, label="QC F1", value="0.2909", detail="缺陷区域指标", accent=COLORS["sky"], detail_size=10)
    add_stat_box(slide, 9.82, 3.22, 2.18, 1.14, label="AUPRC", value="0.8616", detail="缺陷排序能力较好", accent=COLORS["green"], detail_size=10)
    add_panel(
        slide,
        7.40,
        4.58,
        4.60,
        1.52,
        title="这一页最值得强调",
        body="项目已经初步具备“结果可信度筛查”能力。对真实应用来说，这比单纯再提高一个指标更有推进价值。",
        accent=COLORS["navy"],
        fill=COLORS["soft_blue"],
        body_size=14,
    )
    add_panel(slide, 0.72, 5.90, 6.35, 0.92, title="补充解读", body="当前中心区表现明显优于边缘区，说明边缘带仍是最主要的不稳定来源。", accent=COLORS["orange"], fill=COLORS["soft_orange"], title_size=14, body_size=12, title_box_height=0.22, body_top=0.42)
    add_footer(slide, "数据来自成熟度 6 的 QC 评估指标和门控报告。")


def add_pre_real_slide(prs: Presentation, assets: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide)
    add_header(slide, "8. 预真实验证: 已接近真实应用, 但边界清晰", "PRE-REAL VALIDATION", 9)

    slide.shapes.add_picture(str(assets["pre_real"]), Inches(0.72), Inches(1.68), width=Inches(6.9))
    add_stat_box(slide, 7.92, 1.88, 2.10, 1.18, label="BASELINE QC", value="92.5%", detail="基准场景稳定通过", accent=COLORS["green"], detail_size=10)
    add_stat_box(slide, 10.10, 1.88, 2.10, 1.18, label="HI-RES QC", value="0%", detail="256 / 512 明显失稳", accent=COLORS["red"], detail_size=10)
    add_panel(slide, 7.90, 3.24, 4.30, 1.32, title="稳健区域", body="孔径半径 0.95 - 1.05 与小幅 dx 偏差下，QC 通过率基本维持在 0.875 - 1.000。", accent=COLORS["green"], fill=COLORS["soft_green"], body_size=14)
    add_panel(slide, 7.90, 4.78, 4.30, 1.48, title="敏感区域", body="更高分辨率和强噪声会显著拉低 QC 通过率，说明模型仍对强域偏移敏感。", accent=COLORS["red"], fill=COLORS["soft_red"], body_size=14)
    add_footer(slide, "数据来自成熟度 5 数据组 3 和成熟度 6 预真实验证报告。")


def add_summary_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide)
    add_header(slide, "9. 阶段结论与下一步推进", "SUMMARY & NEXT STEP", 10)

    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.74), Inches(1.52), Inches(11.48), Inches(0.84))
    banner.fill.solid()
    banner.fill.fore_color.rgb = rgb(COLORS["navy"])
    banner.line.fill.background()
    add_textbox(slide, 0.98, 1.76, 10.8, 0.28, "当前判断: 路线成立, 证据充分, 边界清楚, 值得继续推进真实闭环。", size=23, color="FFFFFF", bold=True)

    add_panel(
        slide,
        0.78,
        2.78,
        5.45,
        2.58,
        title="已经证明的事情",
        body="1. 方法原型已搭建，且并非停留在概念层面。\n2. 消融实验已经锁定当前优选路径。\n3. OOD 与 QC 结果说明该路线具备可解释的有效性。\n4. 预真实验证表明这条路有现实推进价值。",
        accent=COLORS["blue"],
        body_size=14,
    )
    add_panel(
        slide,
        6.66,
        2.78,
        5.52,
        2.58,
        title="下一阶段最该做的事",
        body="1. 先完成 clean repeat 与 paired defect 的 pilot 采集。\n2. 跑通 qDIC / WLI 最小真实验证闭环。\n3. 用真实样本重新校准 QC 阈值。\n4. 在真实闭环基础上凝练论文初稿和专利技术点。",
        accent=COLORS["green"],
        body_size=14,
    )

    add_panel(slide, 0.78, 5.62, 3.62, 0.96, title="Route works", body="但当前仍不是结题式展示。", accent=COLORS["blue"], fill=COLORS["soft_blue"], title_size=16, body_size=12, title_box_height=0.22, body_top=0.42)
    add_panel(slide, 4.66, 5.62, 3.62, 0.96, title="Evidence chain complete", body="答辩时按证据链逐页推进。", accent=COLORS["teal"], fill=COLORS["card"], title_size=16, body_size=12, title_box_height=0.22, body_top=0.42)
    add_panel(slide, 8.54, 5.62, 3.62, 0.96, title="Real loop next", body="后续重点转向真实闭环。", accent=COLORS["green"], fill=COLORS["soft_green"], title_size=16, body_size=12, title_box_height=0.22, body_top=0.42)
    add_footer(slide, "这版 PPT 以“中期成果型答辩”而非“结题成果型答辩”为定位。")


def build_presentation(assets: dict[str, Path]) -> Presentation:
    prs = new_presentation()
    add_cover_slide(prs, assets)
    add_problem_slide(prs)
    add_route_slide(prs)
    add_method_slide(prs)
    add_ablation_slide(prs, assets)
    add_visual_slide(prs, assets)
    add_ood_slide(prs, assets)
    add_qc_slide(prs, assets)
    add_pre_real_slide(prs, assets)
    add_summary_slide(prs)
    return prs


def validate_presentation(prs: Presentation) -> None:
    max_x = prs.slide_width
    max_y = prs.slide_height
    issues: list[str] = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape_idx, shape in enumerate(slide.shapes, 1):
            if shape.width <= 0 or shape.height <= 0:
                issues.append(f"slide {slide_idx} shape {shape_idx}: non-positive size")
            if shape.left < 0 or shape.top < 0:
                issues.append(f"slide {slide_idx} shape {shape_idx}: negative position")
            if shape.left + shape.width > max_x + 5:
                issues.append(f"slide {slide_idx} shape {shape_idx}: overflow right")
            if shape.top + shape.height > max_y + 5:
                issues.append(f"slide {slide_idx} shape {shape_idx}: overflow bottom")
    if issues:
        raise RuntimeError("Layout validation failed:\n" + "\n".join(issues))


def export_pdf_chunked(pptx_path: Path, pdf_path: Path, chunk_size: int = 3) -> None:
    probe = SpirePresentation()
    probe.LoadFromFile(str(pptx_path))
    slide_count = probe.Slides.Count
    probe.Dispose()

    with tempfile.TemporaryDirectory(prefix="ppt_pdf_chunks_") as tmp:
        chunk_paths: list[Path] = []
        for start in range(0, slide_count, chunk_size):
            end = min(start + chunk_size, slide_count)
            chunk_pdf = Path(tmp) / f"chunk_{start + 1:02d}_{end:02d}.pdf"
            presentation = SpirePresentation()
            presentation.LoadFromFile(str(pptx_path))
            for idx in range(presentation.Slides.Count - 1, -1, -1):
                if idx < start or idx >= end:
                    presentation.Slides.RemoveAt(idx)
            presentation.SaveToFile(str(chunk_pdf), FileFormat.PDF)
            presentation.Dispose()
            chunk_paths.append(chunk_pdf)

        merged = fitz.open()
        try:
            for chunk_pdf in chunk_paths:
                doc = fitz.open(str(chunk_pdf))
                try:
                    merged.insert_pdf(doc)
                finally:
                    doc.close()
            merged.save(str(pdf_path))
        finally:
            merged.close()


def render_pdf_pages(pdf_path: Path, preview_dir: Path, scale: float = 2.0) -> list[Path]:
    if preview_dir.exists():
        shutil.rmtree(preview_dir)
    preview_dir.mkdir(parents=True, exist_ok=True)

    doc = fitz.open(str(pdf_path))
    rendered: list[Path] = []
    try:
        matrix = fitz.Matrix(scale, scale)
        for page_index, page in enumerate(doc, 1):
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            out_path = preview_dir / f"page_{page_index:02d}.png"
            pix.save(str(out_path))
            rendered.append(out_path)
    finally:
        doc.close()
    return rendered


def build_preview_overview(image_paths: list[Path], dst: Path, columns: int = 2, card_bg: str = "#E2E8F0") -> None:
    if not image_paths:
        raise RuntimeError("No preview images available for overview rendering")

    images = [Image.open(path).convert("RGB") for path in image_paths]
    try:
        thumb_width = 900
        gap = 36
        caption_h = 44
        outer = 40
        rows = math.ceil(len(images) / columns)

        resized: list[Image.Image] = []
        max_h = 0
        for img in images:
            scale = thumb_width / img.width
            target_h = int(img.height * scale)
            resized_img = img.resize((thumb_width, target_h), Image.LANCZOS)
            resized.append(resized_img)
            max_h = max(max_h, target_h)

        canvas_w = outer * 2 + columns * thumb_width + (columns - 1) * gap
        cell_h = max_h + caption_h
        canvas_h = outer * 2 + rows * cell_h + (rows - 1) * gap
        canvas = Image.new("RGB", (canvas_w, canvas_h), "white")

        from PIL import ImageDraw

        draw = ImageDraw.Draw(canvas)
        for idx, img in enumerate(resized):
            row = idx // columns
            col = idx % columns
            x = outer + col * (thumb_width + gap)
            y = outer + row * (cell_h + gap)
            draw.rounded_rectangle((x - 10, y - 10, x + thumb_width + 10, y + img.height + caption_h), radius=18, fill=card_bg)
            canvas.paste(img, (x, y))
            draw.text((x + 6, y + img.height + 10), f"Slide {idx + 1}", fill="#334155")

        canvas.save(dst)
    finally:
        for img in images:
            img.close()


def main() -> None:
    assets = build_assets()
    prs = build_presentation(assets)
    validate_presentation(prs)
    prs.save(OUTPUT_PPTX)
    export_pdf_chunked(OUTPUT_PPTX, OUTPUT_PDF)
    preview_pages = render_pdf_pages(OUTPUT_PDF, PREVIEW_DIR, scale=2.0)
    build_preview_overview(preview_pages, PREVIEW_OVERVIEW)
    print(f"Saved PPTX to: {OUTPUT_PPTX}")
    print(f"Saved PDF to: {OUTPUT_PDF}")
    print(f"Preview pages: {PREVIEW_DIR}")
    print(f"Preview overview: {PREVIEW_OVERVIEW}")
    print(f"Assets directory: {ASSETS_DIR}")


if __name__ == "__main__":
    main()
